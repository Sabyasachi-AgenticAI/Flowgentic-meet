from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
# Use a blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Background color (Dark theme)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42) # slate-900

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Flowgentic Meet: The Multi-Agent AI Conferencing Platform"
p.font.bold = True
p.font.size = Pt(30)
p.font.color.rgb = RGBColor(255, 255, 255)

# Add Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
p2 = subtitle_frame.paragraphs[0]
p2.text = "Standalone video conferencing or seamlessly plugged into Zoom, Google Meet, and Teams."
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(148, 163, 184) # slate-400

# Add Main Content - Left Column (The Platform)
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.2), Inches(4.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p_l1 = left_frame.add_paragraph()
p_l1.text = "Our Core Differentiator"
p_l1.font.bold = True
p_l1.font.size = Pt(20)
p_l1.font.color.rgb = RGBColor(56, 189, 248) # sky-400

p_l2 = left_frame.add_paragraph()
p_l2.text = "A Multi-Agent, Multimodal, and Multilingual platform that brings specialized AI intelligence directly into human conversations. Stop relying on post-meeting notes; get answers live."
p_l2.font.size = Pt(14)
p_l2.font.color.rgb = RGBColor(255, 255, 255)
p_l2.space_before = Pt(10)

p_l3 = left_frame.add_paragraph()
p_l3.text = "Data-Driven Decisions"
p_l3.font.bold = True
p_l3.font.size = Pt(20)
p_l3.font.color.rgb = RGBColor(56, 189, 248)
p_l3.space_before = Pt(25)

p_l4 = left_frame.add_paragraph()
p_l4.text = "Agents operate securely in the background with real-time access to your Enterprise Data and External Data sources. This empowers humans to make concise, confident decisions instantly during the meeting."
p_l4.font.size = Pt(14)
p_l4.font.color.rgb = RGBColor(255, 255, 255)
p_l4.space_before = Pt(10)

# Add Main Content - Right Column (The Agents)
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.3), Inches(4.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p_r1 = right_frame.add_paragraph()
p_r1.text = "Specialized AI Personas"
p_r1.font.bold = True
p_r1.font.size = Pt(20)
p_r1.font.color.rgb = RGBColor(167, 139, 250) # violet-400

p_r2 = right_frame.add_paragraph()
p_r2.text = "Summon an individual specialist or an entire 'Agent Swarm' for collaborative advice and suggestions:"
p_r2.font.size = Pt(14)
p_r2.font.color.rgb = RGBColor(255, 255, 255)
p_r2.space_before = Pt(10)

bullet_points = [
    "Project Manager: Timelines, resource allocation, and tracking.",
    "Tax Specialist: Financial structuring and compliance implications.",
    "Compliance Specialist: Real-time regulatory checks.",
    "Senior Tech Lead: System architecture and technical feasibility.",
    "Sales Head: Pitching products, pipeline analysis, and closing strategies."
]

for pt in bullet_points:
    p_bullet = right_frame.add_paragraph()
    p_bullet.text = "• " + pt
    p_bullet.font.size = Pt(13)
    p_bullet.font.color.rgb = RGBColor(226, 232, 240) # slate-200
    p_bullet.space_before = Pt(8)

# Save
prs.save('Flowgentic_Meet_OnePager.pptx')
print("Presentation saved as Flowgentic_Meet_OnePager.pptx")
