import asyncio
import logging
import os
import urllib.parse
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
from livekit import rtc

logger = logging.getLogger("agent")

class SlidePresenter:
    """Manages publishing a local video track as a screen share, rendering slides from a local PPTX file, or streaming a live Playwright browser view."""

    def __init__(self, room) -> None:
        self.room = room
        self.width = 1280
        self.height = 720
        self.source = rtc.VideoSource(self.width, self.height)
        self.track = rtc.LocalVideoTrack.create_video_track("agent-screen-share", self.source)
        self.publication = None
        self.current_slide = 1
        self.running = False
        self.publishing = False
        self.slides_data = []
        
        # Dual presentation mode: 'slides' or 'browser'
        self.mode = "slides"
        self.playwright = None
        self.browser = None
        self.browser_context = None
        self.page = None
        
        # Load premium Windows system font if available, fallback to default
        self.font_title = None
        self.font_header = None
        self.font_body = None
        self.font_small = None
        
        font_paths = [
            "C:\\Windows\\Fonts\\segoeui.ttf",
            "C:\\Windows\\Fonts\\arial.ttf",
            "C:\\Windows\\Fonts\\calibri.ttf",
        ]
        for path in font_paths:
            if os.path.exists(path):
                try:
                    self.font_title = ImageFont.truetype(path, 48)
                    self.font_header = ImageFont.truetype(path, 32)
                    self.font_body = ImageFont.truetype(path, 22)
                    self.font_small = ImageFont.truetype(path, 14)
                    break
                except Exception as e:
                    logger.warning(f"Failed to load font {path}: {e}")
                    
        if not self.font_title:
            self.font_title = ImageFont.load_default()
            self.font_header = ImageFont.load_default()
            self.font_body = ImageFont.load_default()
            self.font_small = ImageFont.load_default()

        # Load initial slide data
        self._load_slide_data()

    def _load_slide_data(self):
        filename = "gtm_presentation.pptx"
        if not os.path.exists(filename):
            logger.warning(f"PPTX file {filename} not found, using static fallback.")
            self._use_fallback_data()
            return
            
        try:
            from pptx import Presentation
            prs = Presentation(filename)
            parsed_slides = []
            for slide in prs.slides:
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text
                bullets = []
                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        continue
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                bullets.append((paragraph.level, text))
                parsed_slides.append({"title": title, "bullets": bullets})
            self.slides_data = parsed_slides
            logger.info(f"Successfully loaded {len(self.slides_data)} slides from {filename}")
        except Exception as e:
            logger.error(f"Error loading PPTX in presenter: {e}")
            self._use_fallback_data()

    def _use_fallback_data(self):
        self.slides_data = [
            {
                "title": "GTM Launch Campaign Status",
                "bullets": [
                    (0, "Launch Timeline & Prep"),
                    (1, "Launch Channel: Product Hunt"),
                    (1, "Date: October 14th (six weeks remaining)"),
                    (1, "Campaign Prep Completion: 55%"),
                    (0, "Primary Target Outreach"),
                    (1, "Prospect: Axcelerate pre-launch private beta"),
                    (1, "Outreach Email: partnerships@axcelerate.io (Status: Ready to send)"),
                ]
            },
            {
                "title": "Sprint Backlog Adjustments",
                "bullets": [
                    (0, "SSO Authentication Integration"),
                    (1, "Swapped In (Priority: Urgent / Launch blocker)"),
                    (1, "Status: In Progress (FLO-4)"),
                    (0, "Custom Dashboard Widget"),
                    (1, "Parked/Archived (Priority: Low / Not customer-facing)"),
                    (1, "Status: Moved to backlog"),
                ]
            },
            {
                "title": "Compliance Assessment",
                "bullets": [
                    (0, "GDPR Onboarding Consent Gap"),
                    (1, "Onboarding step two collects user data without consent checkbox."),
                    (1, "Resolution: Blocking ticket raised on Axcelerate beta invite."),
                    (0, "EU AI Act Compliance Review"),
                    (1, "Automated recommendation model needs transparency check."),
                    (1, "Status: Awaiting AI Act news guidelines check."),
                ]
            }
        ]

    async def start(self) -> str:
        """Starts screen sharing and begins the frame loop."""
        if self.running:
            return "Screen share is already running."
        if self.publishing:
            return "Screen share publication is already in progress."
            
        self.publishing = True
        logger.info("Starting screen share track...")
        options = rtc.TrackPublishOptions()
        options.source = rtc.TrackSource.SOURCE_SCREENSHARE
        
        try:
            self.publication = await self.room.local_participant.publish_track(self.track, options)
            self.running = True
            asyncio.create_task(self._frame_loop())
            return "Screen share started successfully."
        except Exception as e:
            logger.error(f"Failed to start screen share track: {e}")
            return f"Error starting screen share: {e}"
        finally:
            self.publishing = False

    async def stop(self) -> str:
        """Stops screen sharing and unpublishes the track."""
        if not self.running:
            return "Screen share is not running."
            
        self.running = False
        logger.info("Stopping screen share track...")
        try:
            await self.room.local_participant.unpublish_track(self.publication.sid)
            self.publication = None
            return "Screen share stopped successfully."
        except Exception as e:
            logger.error(f"Failed to stop screen share track: {e}")
            return f"Error stopping screen share: {e}"

    def show_slide(self, slide_num: int):
        """Swaps to a specific slide number and switches mode to slides."""
        self._load_slide_data() # Live reload from file in case it was edited
        self.mode = "slides"
        self.current_slide = max(1, min(slide_num, len(self.slides_data)))
        logger.info(f"Switched to slide {self.current_slide}")

    def next_slide(self):
        """Advances to the next slide."""
        self.show_slide(self.current_slide + 1)

    async def start_browser_session(self) -> str:
        """Launches Playwright headless Chromium browser if not already running."""
        if self.browser:
            return "Browser session is already running."
            
        logger.info("Launching headless browser...")
        try:
            from playwright.async_api import async_playwright
            self.playwright = await async_playwright().start()
            
            # Try to launch with system Edge or Chrome first to avoid download issues
            try:
                logger.info("Trying to launch playwright with msedge channel...")
                self.browser = await self.playwright.chromium.launch(headless=True, channel="msedge")
            except Exception as e_edge:
                logger.warning(f"Failed to launch with msedge channel: {e_edge}. Trying chrome...")
                try:
                    self.browser = await self.playwright.chromium.launch(headless=True, channel="chrome")
                except Exception as e_chrome:
                    logger.warning(f"Failed to launch with chrome channel: {e_chrome}. Falling back to default chromium...")
                    self.browser = await self.playwright.chromium.launch(headless=True)

            state_path = "linear_state.json"
            storage_state = state_path if os.path.exists(state_path) else None
            if storage_state:
                logger.info(f"Loading Playwright storage state from: {state_path}")
                
            self.browser_context = await self.browser.new_context(
                storage_state=storage_state,
                viewport={"width": self.width, "height": self.height},
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
                locale="en-US",
                timezone_id="America/New_York"
            )
            self.page = await self.browser_context.new_page()
            
            # Mask webdriver to reduce automated bot flagging
            await self.page.add_init_script(
                "Object.defineProperty(navigator, 'webdriver', {get: () => undefined})"
            )
            
            self.mode = "browser"
            logger.info("Browser session started successfully.")
            return "Browser session started successfully."
        except Exception as e:
            logger.error(f"Failed to start browser session: {e}")
            return f"Error starting browser session: {e}"

    async def browse_to(self, query_or_url: str) -> str:
        """Navigates browser to URL or runs a search."""
        if not self.page:
            return "Error: Browser session is not active."
            
        url = query_or_url.strip()
        if not (url.startswith("http://") or url.startswith("https://")):
            escaped_query = urllib.parse.quote(query_or_url)
            # Use DuckDuckGo to bypass bot detection/CAPTCHA checks
            url = f"https://duckduckgo.com/?q={escaped_query}"
            
        logger.info(f"Navigating browser to: {url}")
        try:
            # We wait up to 12s for load, capturing page rendering state
            await self.page.goto(url, timeout=12000, wait_until="load")
            return f"Navigated successfully to: {url}"
        except Exception as e:
            logger.error(f"Failed navigating to {url}: {e}")
            return f"Navigation completed with warning: {e}"

    async def scroll_browser(self, direction: str) -> str:
        """Scrolls the page view up or down."""
        if not self.page:
            return "Error: Browser session is not active."
            
        try:
            if direction.lower() == "down":
                await self.page.evaluate("window.scrollBy(0, 400);")
                return "Scrolled down successfully."
            elif direction.lower() == "up":
                await self.page.evaluate("window.scrollBy(0, -400);")
                return "Scrolled up successfully."
            else:
                return f"Invalid scroll direction: {direction}."
        except Exception as e:
            logger.error(f"Failed to scroll page: {e}")
            return f"Error scrolling page: {e}"

    async def click_all_issues(self) -> str:
        """Clicks the 'All issues' tab on the Linear dashboard if available."""
        if not self.page:
            return "Error: Browser session is not active."
        try:
            # Give the page a moment to load and render elements
            await asyncio.sleep(2)
            # Try different locator strategies for the "All issues" tab
            locators = [
                self.page.get_by_text("All issues", exact=True),
                self.page.get_by_text("All Issues", exact=True),
                self.page.locator('button:has-text("All issues")'),
                self.page.locator('div:has-text("All issues")'),
                self.page.locator('a:has-text("All issues")'),
                self.page.locator('[data-testid="all-issues-tab"]')
            ]
            for loc in locators:
                try:
                    count = await loc.count()
                    for i in range(count):
                        element = loc.nth(i)
                        if await element.is_visible():
                            await element.click()
                            logger.info("Clicked 'All issues' tab successfully.")
                            return "Clicked 'All issues' tab successfully."
                except Exception as ex:
                    logger.debug(f"Locator check failed: {ex}")
            
            # Fallback to JS evaluation to scan the DOM and click
            clicked = await self.page.evaluate('''() => {
                const elements = Array.from(document.querySelectorAll('*'));
                for (const el of elements) {
                    if (el.textContent && el.textContent.trim() === 'All issues') {
                        el.click();
                        return true;
                    }
                }
                for (const el of elements) {
                    if (el.textContent && el.textContent.trim().toLowerCase() === 'all issues') {
                        el.click();
                        return true;
                    }
                }
                return false;
            }''')
            if clicked:
                logger.info("Clicked 'All issues' via JS evaluation.")
                return "Clicked 'All issues' tab via JS."
            
            return "Could not find 'All issues' tab to click."
        except Exception as e:
            logger.error(f"Failed during click_all_issues: {e}")
            return f"Error clicking 'All issues': {e}"

    async def stop_browser_session(self) -> str:
        """Closes browser session."""
        if not self.browser:
            return "Browser session is not active."
            
        logger.info("Closing browser session...")
        try:
            if self.page:
                await self.page.close()
            if self.browser_context:
                await self.browser_context.close()
            if self.browser:
                await self.browser.close()
            if self.playwright:
                await self.playwright.stop()
        except Exception as e:
            logger.error(f"Error closing browser: {e}")
        finally:
            self.page = None
            self.browser_context = None
            self.browser = None
            self.playwright = None
            self.mode = "slides"
            
        return "Browser session closed."

    async def _frame_loop(self):
        """Sends frames continuously at ~5 FPS to keep WebRTC stream active with low bandwidth."""
        while self.running:
            try:
                img = await self._render_current_slide()
                rgba_data = img.tobytes()
                rgba_frame = rtc.VideoFrame(self.width, self.height, rtc.VideoBufferType.RGBA, rgba_data)
                # Convert to I420 (YUV) color format to prevent green screen issues in WebRTC
                i420_frame = rgba_frame.convert(rtc.VideoBufferType.I420)
                self.source.capture_frame(i420_frame)
            except Exception as e:
                logger.error(f"Error in slide frame loop: {e}")
            await asyncio.sleep(0.2) # 5 fps is plenty for slides

    async def _render_current_slide(self) -> Image.Image:
        """Renders the Pillow image based on mode (slides or browser)."""
        # If in browser mode, take screenshot of the active browser page
        if self.mode == "browser" and self.page:
            try:
                screenshot_bytes = await self.page.screenshot(type="png")
                img = Image.open(BytesIO(screenshot_bytes))
                if img.width != self.width or img.height != self.height:
                    img = img.resize((self.width, self.height), Image.Resampling.LANCZOS)
                return img.convert("RGBA")
            except Exception as e:
                logger.error(f"Failed to capture browser screenshot: {e}")
                # Fallback to drawing standard slides below if screenshot fails
        
        # Create dark blue/violet gradient image background
        img = Image.new("RGBA", (self.width, self.height), (11, 13, 25, 255))
        draw = ImageDraw.Draw(img)
        
        # Subtle glassmorphic background accent glows
        draw.ellipse([800, -200, 1400, 400], fill=(139, 92, 246, 25))  # Purple glow
        draw.ellipse([-200, 400, 400, 900], fill=(6, 182, 212, 15))    # Cyan glow

        # Top Navigation Bar
        draw.rectangle([0, 0, self.width, 70], fill=(20, 24, 45, 255))
        draw.line([0, 70, self.width, 70], fill=(45, 50, 80, 255), width=1)
        
        # Title of the room in nav bar
        draw.text((40, 20), "FLOWSYNC WAR ROOM", fill=(255, 255, 255, 255), font=self.font_small)
        draw.text((1120, 20), "NexaCore Inc.", fill=(139, 92, 246, 255), font=self.font_small)

        # Bottom Bar info
        draw.rectangle([0, self.height-40, self.width, self.height], fill=(20, 24, 45, 255))
        draw.line([0, self.height-40, self.width, self.height-40], fill=(45, 50, 80, 255), width=1)
        draw.text((40, self.height-30), f"Slide {self.current_slide} of {len(self.slides_data)}", fill=(156, 163, 175, 255), font=self.font_small)
        draw.text((1100, self.height-30), "AI Assistant Presenting", fill=(6, 182, 212, 255), font=self.font_small)

        self._render_dynamic_slide(draw, self.current_slide - 1)
            
        return img

    def _render_dynamic_slide(self, draw: ImageDraw.ImageDraw, slide_idx: int):
        if slide_idx < 0 or slide_idx >= len(self.slides_data):
            return
            
        slide = self.slides_data[slide_idx]
        title = slide["title"]
        bullets = slide["bullets"]
        
        # Slide Title
        draw.text((80, 110), title, fill=(255, 255, 255, 255), font=self.font_title)
        
        # Group bullets by level-0 items to form cards dynamically
        cards = []
        current_card = None
        for level, text in bullets:
            if level == 0:
                if current_card:
                    cards.append(current_card)
                current_card = {"header": text, "items": []}
            else:
                if not current_card:
                    current_card = {"header": "Details", "items": []}
                current_card["items"].append(text)
        if current_card:
            cards.append(current_card)
            
        # Draw cards based on count
        if not cards:
            return
            
        if len(cards) == 1:
            c = cards[0]
            card_rect = [80, 220, 1200, 600]
            outline_color = (245, 158, 11, 255) if "compliance" in title.lower() else (139, 92, 246, 255)
            draw.rectangle(card_rect, fill=(20, 24, 45, 200), outline=outline_color, width=2)
            draw.text((110, 250), c["header"], fill=outline_color, font=self.font_header)
            
            y = 320
            for item in c["items"]:
                draw.text((110, y), f"• {item}", fill=(255, 255, 255, 255), font=self.font_body)
                y += 50
                
        else:
            c1 = cards[0]
            card1_rect = [80, 220, 600, 600]
            col1 = (16, 185, 129, 255) if "sso" in c1["header"].lower() or "auth" in c1["header"].lower() else (139, 92, 246, 255)
            draw.rectangle(card1_rect, fill=(20, 24, 45, 200), outline=col1, width=2)
            draw.text((110, 250), c1["header"], fill=col1, font=self.font_header)
            
            y = 320
            for item in c1["items"]:
                draw.text((110, y), f"• {item}", fill=(255, 255, 255, 255), font=self.font_body)
                y += 50
                
            c2 = cards[1]
            card2_rect = [680, 220, 1200, 600]
            col2 = (239, 68, 68, 255) if "parked" in c2["header"].lower() or "widget" in c2["header"].lower() or "archived" in c2["header"].lower() else (6, 182, 212, 255)
            draw.rectangle(card2_rect, fill=(20, 24, 45, 200), outline=col2, width=2)
            draw.text((710, 250), c2["header"], fill=col2, font=self.font_header)
            
            y = 320
            for item in c2["items"]:
                draw.text((710, y), f"• {item}", fill=(255, 255, 255, 255), font=self.font_body)
                y += 50
