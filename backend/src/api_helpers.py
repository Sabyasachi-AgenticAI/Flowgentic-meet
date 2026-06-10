import base64
import json
import logging
import os
import smtplib
import time
import urllib.request
import uuid
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

logger = logging.getLogger("agent")

# Linear API configurations
LINEAR_TEAM_ID = "95542484-a50b-4587-999a-a682c3fe7cad"
LINEAR_STATES = {
    "Backlog": "0f45246b-dbeb-4414-8a50-c0c0bd881157",
    "Todo": "3862f813-640e-4a02-bcd5-fdc99c564940",
    "In Progress": "2e31fea4-5438-464c-ab7d-853aabf3d8ed",
    "Done": "3d4308f2-b85f-42de-9241-eec2379c11f5",
}


def call_linear_api(query: str, variables: dict | None = None) -> dict | None:
    api_key = os.getenv("LINEAR_API_KEY")
    if not api_key:
        logger.error("LINEAR_API_KEY not set")
        return None
    req_data = json.dumps({"query": query, "variables": variables or {}}).encode(
        "utf-8"
    )
    req = urllib.request.Request(
        "https://api.linear.app/graphql",
        data=req_data,
        headers={"Content-Type": "application/json", "Authorization": api_key},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req) as response:
            res_data = response.read().decode("utf-8")
            return json.loads(res_data)
    except Exception as e:
        logger.error(f"Error calling Linear API: {e}")
        return None


def call_devops_api(
    endpoint_suffix: str, method: str = "GET", body: list | dict | None = None
) -> dict | None:
    pat = os.getenv("AZURE_DEVOPS_PAT")
    org = os.getenv("AZURE_DEVOPS_ORG")
    project = os.getenv("AZURE_DEVOPS_PROJECT")
    if not pat or not org or not project:
        logger.error("Azure DevOps credentials not fully set")
        return None

    org_url = org.rstrip("/")
    url = f"{org_url}/{project}/_apis/{endpoint_suffix.lstrip('/')}"

    auth_val = base64.b64encode(f":{pat}".encode()).decode("utf-8")
    headers = {
        "Content-Type": "application/json-patch+json"
        if method in ("POST", "PATCH")
        else "application/json",
        "Authorization": f"Basic {auth_val}",
    }

    req_data = json.dumps(body).encode("utf-8") if body is not None else None
    req = urllib.request.Request(url, data=req_data, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req) as response:
            res_data = response.read().decode("utf-8")
            return json.loads(res_data)
    except Exception as e:
        logger.error(f"Error calling Azure DevOps API: {e}")
        return None


def call_n8n_webhook(url: str, body: dict) -> dict | None:
    if not url:
        return None
    req_data = json.dumps(body).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=req_data,
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    try:
        with urllib.request.urlopen(req) as response:
            res_data = response.read().decode("utf-8")
            if not res_data:
                return {"success": True}
            try:
                return json.loads(res_data)
            except json.JSONDecodeError:
                return {"success": True, "raw": res_data}
    except Exception as e:
        logger.error(f"Error calling n8n webhook {url}: {e}")
        return None


def find_devops_work_item_by_linear_id(linear_id: str) -> int | None:
    query = f"SELECT [System.Id] FROM WorkItems WHERE [System.Title] CONTAINS '[{linear_id}]'"
    res = call_devops_api(
        "wit/wiql?api-version=6.0", method="POST", body={"query": query}
    )
    if res and "workItems" in res and res["workItems"]:
        return res["workItems"][0]["id"]
    return None


def resolve_linear_issue(ticket_id: str | int) -> dict | None:
    ticket_str = str(ticket_id).strip()
    if not ticket_str.upper().startswith("FLO-"):
        issue_number_str = ticket_str
    else:
        issue_number_str = ticket_str[4:]

    try:
        issue_number = int(issue_number_str)
    except ValueError:
        logger.error(f"Invalid ticket number format: {ticket_str}")
        return None

    query = """
    query GetIssue($teamId: String!, $number: Float!) {
      issues(filter: { team: { id: { eq: $teamId } }, number: { eq: $number } }) {
        nodes {
          id
          identifier
          title
          description
          priority
          state {
            id
            name
          }
        }
      }
    }
    """
    res = call_linear_api(query, {"teamId": LINEAR_TEAM_ID, "number": issue_number})
    if res and "data" in res and res["data"]["issues"]["nodes"]:
        return res["data"]["issues"]["nodes"][0]
    return None


# Mock Ticket Database shared across all agents in the session
TICKETS = {
    1: {
        "id": 1,
        "title": "Implement authentication using Clerk",
        "description": "Add login/signup routes and secure meeting rooms.",
        "priority": "High",
        "state": "Backlog",
        "archived": False,
    },
    2: {
        "id": 2,
        "title": "Configure LiveKit room options",
        "description": "Optimize noise cancellation and turn detection.",
        "priority": "Medium",
        "state": "In Progress",
        "archived": False,
    },
    3: {
        "id": 3,
        "title": "Design dark luxury landing page",
        "description": "Create glassmorphic UI layout with rich animations.",
        "priority": "High",
        "state": "Done",
        "archived": False,
    },
}
NEXT_TICKET_ID = 4


def ensure_default_pptx_exists():
    filename = "gtm_presentation.pptx"
    if os.path.exists(filename):
        return
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: GTM Launch Campaign Status
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide1 = prs.slides.add_slide(slide_layout)
        slide1.shapes.title.text = "GTM Launch Campaign Status"
        body_shape = slide1.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Launch Timeline & Prep"
        
        p = tf.add_paragraph()
        p.text = "Launch Channel: Product Hunt"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Date: October fourteenth (six weeks remaining)"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Campaign Prep Completion: fifty-five percent"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Primary Target Outreach"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "Prospect: Axcelerate pre-launch private beta"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Outreach Email: partnerships@axcelerate.io (Status: Ready to send)"
        p.level = 1
        
        # Slide 2: Sprint Backlog Adjustments
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Sprint Backlog Adjustments"
        body_shape2 = slide2.shapes.placeholders[1]
        tf2 = body_shape2.text_frame
        tf2.text = "SSO Authentication Integration"
        
        p = tf2.add_paragraph()
        p.text = "Swapped In (Priority: Urgent / Launch blocker)"
        p.level = 1
        p = tf2.add_paragraph()
        p.text = "Status: In Progress (FLO-four)"
        p.level = 1
        
        p = tf2.add_paragraph()
        p.text = "Custom Dashboard Widget"
        p.level = 0
        p = tf2.add_paragraph()
        p.text = "Parked/Archived (Priority: Low / Not customer-facing)"
        p.level = 1
        p = tf2.add_paragraph()
        p.text = "Status: Moved to backlog"
        p.level = 1
        
        # Slide 3: Compliance Assessment
        slide3 = prs.slides.add_slide(slide_layout)
        slide3.shapes.title.text = "Compliance Assessment"
        body_shape3 = slide3.shapes.placeholders[1]
        tf3 = body_shape3.text_frame
        tf3.text = "GDPR Onboarding Consent Gap"
        
        p = tf3.add_paragraph()
        p.text = "Onboarding step two collects user data without consent checkbox."
        p.level = 1
        p = tf3.add_paragraph()
        p.text = "Resolution: Blocking ticket raised on Axcelerate beta invite."
        p.level = 1
        
        p = tf3.add_paragraph()
        p.text = "EU AI Act Compliance Review"
        p.level = 0
        p = tf3.add_paragraph()
        p.text = "Automated recommendation model needs transparency check."
        p.level = 1
        p = tf3.add_paragraph()
        p.text = "Status: Awaiting AI Act news guidelines check."
        p.level = 1
        
        prs.save(filename)
        logger.info(f"Created default PowerPoint file: {filename}")
    except Exception as e:
        logger.error(f"Failed to create default PPTX file: {e}")


def get_pptx_summary() -> str:
    filename = "gtm_presentation.pptx"
    if not os.path.exists(filename):
        return "No presentation file found."
        
    try:
        from pptx import Presentation
        prs = Presentation(filename)
        summary_lines = []
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            summary_lines.append(f"Slide {slide_num}: {title}")
            
            # Look for all shapes that have text frames
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            indent = "  " * (paragraph.level + 1)
                            summary_lines.append(f"{indent}- {text}")
                            
        return "\n".join(summary_lines)
    except Exception as e:
        logger.error(f"Error parsing PPTX presentation: {e}")
        return f"Error reading presentation: {e}"


def send_outlook_smtp_email(recipient: str, subject: str, body: str) -> bool:
    server_addr = os.getenv("SMTP_SERVER", "smtp.office365.com")
    try:
        port = int(os.getenv("SMTP_PORT", "587"))
    except (TypeError, ValueError):
        port = 587
        
    sender = os.getenv("OUTLOOK_EMAIL")
    password = os.getenv("OUTLOOK_PASSWORD")
    
    if not sender or not password:
        logger.error("SMTP credentials not fully set in environment")
        return False
        
    msg = MIMEMultipart()
    msg["From"] = sender
    msg["To"] = recipient
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain"))
    
    try:
        with smtplib.SMTP(server_addr, port) as server:
            server.starttls()
            server.login(sender, password)
            server.send_message(msg)
        logger.info(f"Outlook SMTP email sent successfully to {recipient}")
        return True
    except Exception as e:
        logger.error(f"Failed to send Outlook SMTP email: {e}")
        return False


def send_teams_webhook_ping(webhook_url: str, text: str) -> bool:
    if not webhook_url:
        logger.error("Teams webhook URL is empty")
        return False
    payload = {"text": text}
    req_data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        webhook_url,
        data=req_data,
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    try:
        with urllib.request.urlopen(req) as response:
            response.read()
        logger.info("Teams webhook ping sent successfully")
        return True
    except Exception as e:
        logger.error(f"Failed to ping Teams webhook: {e}")
        return False


# Generate the default pptx file immediately on module load
ensure_default_pptx_exists()

